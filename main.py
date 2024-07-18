# 진입점
import tkinter as tk
from tkinter import ttk, messagebox, filedialog
from pptx import Presentation
import json
import os

def replace_text_in_shape(shape, replacements):
    if not shape.has_text_frame:
        return

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for placeholder, replacement in replacements.items():
                if f"{placeholder}" in run.text:
                    run.text = run.text.replace(f"{placeholder}", replacement)


def replace_placeholders(ppt_template_path, replacements, output_path):
    prs = Presentation(ppt_template_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                replace_text_in_shape(shape, replacements)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        replace_text_in_shape(cell, replacements)
            if shape.has_text_frame:
                replace_text_in_shape(shape, replacements)

    prs.save(output_path)


class ReplacementManager:
    def __init__(self, root):
        self.root = root
        self.root.title("Replacements Manager")

        self.replacements_file = "replacements.json"
        self.replacements = self.load_replacements()

        self.tree = ttk.Treeview(root, columns=("Key", "Value"), show="headings")
        self.tree.heading("Key", text="Key")
        self.tree.heading("Value", text="Value")
        self.tree.grid(row=0, column=0, columnspan=4)
        self.tree.bind("<<TreeviewSelect>>", self.on_tree_select)

        self.add_key_label = tk.Label(root, text="Key:")
        self.add_key_label.grid(row=1, column=0, sticky=tk.W)
        self.add_key_entry = tk.Entry(root)
        self.add_key_entry.grid(row=1, column=1, sticky=tk.E)

        self.add_value_label = tk.Label(root, text="Value:")
        self.add_value_label.grid(row=2, column=0, sticky=tk.W)
        self.add_value_entry = tk.Entry(root)
        self.add_value_entry.grid(row=2, column=1, sticky=tk.E)

        self.add_button = tk.Button(root, text="Add/Update", command=self.add_or_update)
        self.add_button.grid(row=3, column=0)

        self.delete_button = tk.Button(root, text="Delete", command=self.delete)
        self.delete_button.grid(row=3, column=1)

        self.select_template_button = tk.Button(root, text="Select Template PPT", command=self.select_template)
        self.select_template_button.grid(row=3, column=2)

        self.save_button = tk.Button(root, text="Save Output PPT", command=self.save_output)
        self.save_button.grid(row=3, column=3)

        self.template_path_label = tk.Label(root, text="No template selected", fg="blue")
        self.template_path_label.grid(row=4, column=0, columnspan=4)

        self.output_path_label = tk.Label(root, text="No output file selected", fg="green")
        self.output_path_label.grid(row=5, column=0, columnspan=4)

        self.refresh_button = tk.Button(root, text="Refresh", command=self.refresh_tree)
        self.refresh_button.grid(row=6, column=0, columnspan=4)

        self.ppt_template_path = ""
        self.output_path = ""

        self.refresh_tree()

    def load_replacements(self):
        if os.path.exists(self.replacements_file):
            with open(self.replacements_file, "r", encoding="utf-8") as file:
                return json.load(file)
        return {}

    def save_replacements(self):
        with open(self.replacements_file, "w", encoding="utf-8") as file:
            json.dump(self.replacements, file, indent=4, ensure_ascii=False)

    def add_or_update(self):
        key = self.add_key_entry.get()
        value = self.add_value_entry.get()
        if key and value:
            self.replacements[key] = value
            self.save_replacements()
            self.refresh_tree()
            self.add_key_entry.delete(0, tk.END)
            self.add_value_entry.delete(0, tk.END)
        else:
            messagebox.showwarning("Input Error", "Both key and value must be provided.")

    def delete(self):
        selected_item = self.tree.selection()
        if selected_item:
            key = self.tree.item(selected_item)["values"][0]
            del self.replacements[key]
            self.save_replacements()
            self.refresh_tree()
        else:
            messagebox.showwarning("Selection Error", "No item selected for deletion.")

    def refresh_tree(self):
        for item in self.tree.get_children():
            self.tree.delete(item)
        self.replacements = self.load_replacements()
        for key, value in self.replacements.items():
            self.tree.insert("", "end", values=(key, value))

    def on_tree_select(self, event):
        selected_item = self.tree.selection()
        if selected_item:
            key, value = self.tree.item(selected_item)["values"]
            self.add_key_entry.delete(0, tk.END)
            self.add_value_entry.delete(0, tk.END)
            self.add_key_entry.insert(0, key)
            self.add_value_entry.insert(0, value)

    def select_template(self):
        self.ppt_template_path = filedialog.askopenfilename(
            filetypes=[("PowerPoint files", "*.pptx")])
        if self.ppt_template_path:
            self.template_path_label.config(text=f"Template: {self.ppt_template_path}")

    def save_output(self):
        self.output_path = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint files", "*.pptx")])
        if self.output_path:
            replace_placeholders(self.ppt_template_path, self.replacements, self.output_path)
            self.output_path_label.config(text=f"Output: {self.output_path}")
            messagebox.showinfo("Success", f"Output PPT saved: {self.output_path}")
            
    def save_ppt(self):
        ppt_template_path = "template_vworld.pptx"
        output_path = "[가이드] 브이월드 API 신청_" + self.replacements["name"] + ".pptx"
        replace_placeholders(ppt_template_path, self.replacements, output_path)
        messagebox.showinfo("Success", "PPT saved successfully!")

if __name__ == "__main__":
    replacements = {};
    root = tk.Tk()
    app = ReplacementManager(root)
    root.mainloop()